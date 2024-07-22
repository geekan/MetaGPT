import pdb
from datetime import datetime
from typing import Dict
import os
from pptx import Presentation
import json
# 设置环境变量
import asyncio
from metagpt.logs import logger
from metagpt.roles.role import Role,RoleReactMode
from metagpt.schema import Message
from pptx.util import Pt
from metagpt.utils.file import File
from metagpt.actions import Action
from metagpt.utils.common import OutputParser


class CreatePPTOutline(Action):
    name: str = "CreatePPTOutline",
    language:str = "chinese"

    async def run(self, topic: str, *args, **kwargs) -> Dict:
        COMMON_PROMPT = """
        You are now a seasoned professional in creating PowerPoint presentations.
        We need you to create an outline for a PPT with the topic "{topic}".
        """

        OUTLINE_PROMPT = COMMON_PROMPT + """
        Please provide a specific outline for this PPT, strictly following the following requirements:
        1. The output must be strictly in the specified language, {language}.
        2. Answer strictly in the dictionary format like {{"title": "xxx", "outline": [{{"slide 1": ["point 1", "point 2"]}}, {{"slide 2": ["point 3", "point 4"]}}]}}.
        3. The outline should be as specific and sufficient as possible, with a main slide and sub-points. The sub-points should use {language}.
        4. Do not have extra spaces or line breaks.
        5. Each slide title and point must have practical significance.
        """
        prompt = OUTLINE_PROMPT.format(topic=topic, language=self.language)
        resp = await self._aask(prompt=prompt)
        # 替换全角逗号为半角逗号
        resp = resp.replace('，', ',')
        return OutputParser.extract_struct(resp, "dict")

class CreatePPTContent(Action):
    # def __init__(self, name: str = "CreatePPTContent", outline: str = "", language: str = "chinese", *args, **kwargs):
    #     super().__init__(name, *args, **kwargs)
    #     self.language = language
    #     self.outline = outline
    name: str = "CreatePPTContent"
    outline: str = ""
    language: str = "chinese"

    async def run(self, topic: str, *args, **kwargs) -> str:
        COMMON_PROMPT = """
        You are now a seasoned professional in creating PowerPoint presentations.
        We need you to create content for a PPT with the topic "{topic}".
        """

        output_format = json.dumps({
            "title": "example title",
            "pages": [
                {
                    "title": "title for page 1",
                    "content": [
                        {
                            "title": "title for paragraph 1",
                            "description": ["detail for paragraph 1"],
                        },
                        {
                            "title": "title for paragraph 2",
                            "description": ["detail for paragraph 2"],
                        },
                    ],
                }
            ],
        }, ensure_ascii=False)

        CONTENT_PROMPT = COMMON_PROMPT+"""
        Now I will give you the outline for the topic.
        Please output the detailed content for each slide, ensuring it is well-structured and engaging.

        The outline for the topic is as follows:
        {outline}

        Strictly limit output according to the following requirements:
        1. Follow standard presentation structure and formatting.
        2. Each slide should have a title and detailed content with bullet points or numbered lists.
        3. The output must be strictly in the specified language,{language}.
        4. Do not have redundant output, including concluding remarks.
        5. Strict requirement not to output the topic "{topic}".
        6. According to this JSON format, the output {output_format} can only return JSON, and the JSON should not be wrapped with ```
        """

        prompt = CONTENT_PROMPT.format(
            topic=topic, language=self.language, outline=self.outline,output_format=output_format)
        return await self._aask(prompt=prompt)

class PPTDesigner(Role):

    name: str = "Bob",
    profile: str = "PPTDesigner"
    language: str = "Chinese"
    topic: str = ""
    main_title: str = ""
    total_content: str = ""

    def __init__(self,**kwargs):
        super().__init__(**kwargs)
        self.set_actions([CreatePPTOutline(language=self.language)])
        self._set_react_mode(react_mode=RoleReactMode.BY_ORDER.value)

    async def _think(self) -> None:
        logger.info(self.rc.state)
        if self.rc.todo is None:
            self._set_state(0)
            return

        if self.rc.state + 1 < len(self.states):
            self._set_state(self.rc.state + 1)
        else:
            self.rc.todo = None

    async def _handle_outline(self, outline: Dict) -> Message:
        self.main_title = outline.get("title")
        outline_str = f"{self.main_title}\n"
        self.total_content += f"{self.main_title}"
        actions = list()
        for slide in outline.get("outline"):
            actions.append(CreatePPTContent(
                language=self.language, outline=json.dumps(slide, ensure_ascii=False)))
            key = list(slide.keys())[0]
            outline_str += f"-{key}\n"
            for point in slide[key]:
                outline_str += f"  - {point}\n"
        self.set_actions(actions)
        self.rc.todo = None
        return Message(content=outline_str)

    async def _act(self) -> Message:
        todo = self.rc.todo
        if isinstance(todo, CreatePPTOutline):
            msg = self.rc.memory.get(k=1)[0]
            self.topic = msg.content
            resp = await todo.run(topic=self.topic)
            logger.info(resp)
            return await self._handle_outline(resp)
        resp = await todo.run(topic=self.topic)
        logger.info(resp)
        if self.total_content != "":
            self.total_content += "\n\n\n"
        self.total_content += resp
        return Message(content=resp, role=self.profile)

    async def _react(self) -> Message:
        ppt = None
        first_iteration = True

        while True:
            await self._think()
            if self.rc.todo is None:
                break
            msg = await self._act()
            msg.content = msg.content.replace('，', ',')
            if first_iteration:
                ppt = Presentation()
                slide_layout = ppt.slide_layouts[0]
                slide = ppt.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = msg.content.splitlines()[0]
                subtitle.text = "metaPPT"
                first_iteration = False
            else:
                try:
                    msg_dict = json.loads(msg.content)
                    for page in msg_dict['pages']:
                        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                        slide.shapes.title.text = page['title']
                        text_frame = slide.placeholders[1].text_frame
                        for sub_content in page['content']:
                            # 添加一级内容
                            p = text_frame.add_paragraph()
                            p.text = sub_content['title']
                            p.level = 0
                            p.font.size = Pt(24)

                            # 添加二级内容
                            for description in sub_content['description']:
                                p = text_frame.add_paragraph()
                                p.text = description
                                p.level = 1
                                p.font.size = Pt(18)
                except Exception as e:
                    logger.error(e)

        if ppt:
            ppt.save(f"{self.main_title}.ppt")
        return Message(content=f"PPT {self.main_title} 已生成", role=self.profile)


async def main():
    msg = "华为公司"
    role = PPTDesigner()
    logger.info(msg)
    result = await role.run(msg)
    logger.info(result)

asyncio.run(main())
